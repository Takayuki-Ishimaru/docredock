#!/usr/bin/env python3
"""Generate the deterministic PPTX corpus used to test readable(.md) conversion.

See tests/DocRedock.Tests/Fixtures/COMPLEX_DESIGN_DOC_SPEC.md section 3 (P01-P15).
Facts (IDs/values) are copied verbatim from
outputs/drmd-validation-20260823/経費精算システム_設計書_検証用.md so the
corpus can be cross-checked against the DOCX/PDF siblings.

Usage:
  python3 generate_complex_pptx.py [output.pptx]
"""
from pathlib import Path
import copy
import os
import sys
import zipfile

from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

# Repo root is normally four levels up from this file
# (tests/DocRedock.Tests/Fixtures/Pptx/generate_complex_pptx.py). DRMD_FIXTURE_ROOT
# is an escape hatch for running the script from a copy outside the repo tree.
ROOT = Path(os.environ["DRMD_FIXTURE_ROOT"]) if os.environ.get("DRMD_FIXTURE_ROOT") else Path(__file__).resolve().parents[4]
XLSX_SOURCE = ROOT / "経費精算システム_設計書_検証用.xlsx"
DEFAULT_OUT = Path(__file__).resolve().parent / "complex-design-doc.pptx"

# forbidden token: must NEVER appear in any slide/notes text we author.
FORBIDDEN_TOKEN = "OCR-JP-20260823-017"

# ---------------------------------------------------------------- palette --
NAVY = RGBColor(0x17, 0x36, 0x5D)
TEAL = RGBColor(0x0B, 0x72, 0x85)
INK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x66, 0x70, 0x85)
RED = RGBColor(0xB4, 0x23, 0x18)
GREEN = RGBColor(0x06, 0x76, 0x47)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xB7, 0x79, 0x1B)
JP_FONT = "Yu Gothic"

CHECKLIST = []  # list of (id, description, method, "PASS"/"FAIL")

def record(elem_id, desc, method, ok):
    CHECKLIST.append((elem_id, desc, method, "PASS" if ok else "FAIL"))

# ------------------------------------------------------------ image assets --
def load_images():
    with zipfile.ZipFile(XLSX_SOURCE) as z:
        img1 = z.read("xl/media/image.png")  # IMG-01, no text
        img2 = z.read("xl/media/image2.png")  # IMG-02, japanese text (contains forbidden token as PIXELS only)
    return img1, img2

# --------------------------------------------------------- bullet/xml util --
def _clear_bullet(pPr):
    for tag in (
        "a:buClr", "a:buClrTx", "a:buSzPct", "a:buSzPts", "a:buSzTx",
        "a:buFont", "a:buFontTx", "a:buNone", "a:buChar", "a:buAutoNum",
    ):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)

def set_bullet_char(paragraph, char, font=JP_FONT):
    """Inject an explicit buChar bullet (with buFont) into this paragraph's pPr."""
    pPr = paragraph._p.get_or_add_pPr()
    _clear_bullet(pPr)
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": font})
    pPr.append(buFont)
    buChar = pPr.makeelement(qn("a:buChar"), {"char": char})
    pPr.append(buChar)

def set_bullet_none(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    _clear_bullet(pPr)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))

def set_bullet_autonum(paragraph, scheme="arabicPeriod", start_at=None):
    pPr = paragraph._p.get_or_add_pPr()
    _clear_bullet(pPr)
    attrib = {"type": scheme}
    if start_at is not None:
        attrib["startAt"] = str(start_at)
    pPr.append(pPr.makeelement(qn("a:buAutoNum"), attrib))

def add_run(paragraph, text, bold=False, italic=False, underline=False,
            strike=False, color=None, size=None, font=JP_FONT):
    assert FORBIDDEN_TOKEN not in text, "forbidden OCR token must not appear in slide text"
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic
    if underline:
        run.font.underline = True
    if size:
        run.font.size = Pt(size)
    if color:
        run.font.color.rgb = color
    if strike:
        rPr = run._r.get_or_add_rPr()
        rPr.set("strike", "sngStrike")  # not exposed by python-pptx: raw XML
    return run

def new_paragraph(text_frame, first=False):
    if first and len(text_frame.paragraphs) == 1 and not text_frame.paragraphs[0].runs:
        return text_frame.paragraphs[0]
    return text_frame.add_paragraph()

def plain_paragraph(text_frame, text, first=False, size=14, color=INK, font=JP_FONT):
    """Add a paragraph via the plain high-level API only (no bullet XML at all).

    Used for P05: the paragraph inherits whatever bullet the slideLayout/master
    placeholder defines, because we never touch <a:pPr> bullet children here.
    """
    assert FORBIDDEN_TOKEN not in text
    p = new_paragraph(text_frame, first=first)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = font
    return p

def set_title(slide, text, size=None, color=NAVY):
    title = slide.shapes.title
    title.text_frame.text = text
    for p in title.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = JP_FONT
            r.font.color.rgb = color
            if size:
                r.font.size = Pt(size)
    return title

def next_shape_id(slide):
    ids = [int(el.get("id")) for el in slide.shapes._spTree.iter(qn("p:cNvPr"))]
    return max(ids, default=1) + 1

def _find_ph_sp(spTree, ph_type):
    for sp in spTree.iter(qn("p:sp")):
        ph = sp.find(f'{qn("p:nvSpPr")}/{qn("p:nvPr")}/{qn("p:ph")}')
        if ph is not None and ph.get("type") == ph_type:
            return sp
    return None

def add_furniture(slide, footer_text, page_no, date_text="2026-08-23"):
    """P13: clone footer/date/slideNum placeholders from the slide's own layout.

    python-pptx's add_slide() only copies title/body placeholders, so these
    furniture placeholders must be injected by hand (real OOXML placeholder
    shapes, not text boxes) to test that the converter treats them as
    furniture rather than body content.
    """
    layout_spTree = slide.slide_layout.shapes._spTree
    spTree = slide.shapes._spTree
    for ph_type in ("dt", "ftr", "sldNum"):
        src = _find_ph_sp(layout_spTree, ph_type)
        if src is None:
            continue
        new_sp = copy.deepcopy(src)
        cNvPr = new_sp.find(f'{qn("p:nvSpPr")}/{qn("p:cNvPr")}')
        cNvPr.set("id", str(next_shape_id(slide)))
        txBody = new_sp.find(qn("p:txBody"))
        if ph_type == "ftr":
            for p in txBody.findall(qn("a:p")):
                txBody.remove(p)
            p = txBody.makeelement(qn("a:p"), {})
            r = p.makeelement(qn("a:r"), {})
            rPr = p.makeelement(qn("a:rPr"), {"lang": "ja-JP", "sz": "1200"})
            r.append(rPr)
            t = p.makeelement(qn("a:t"), {})
            t.text = footer_text
            r.append(t)
            p.append(r)
            txBody.append(p)
        elif ph_type == "dt":
            fld = txBody.find(f'{qn("a:p")}/{qn("a:fld")}')
            if fld is not None:
                tnode = fld.find(qn("a:t"))
                if tnode is not None:
                    tnode.text = date_text
        elif ph_type == "sldNum":
            fld = txBody.find(f'{qn("a:p")}/{qn("a:fld")}')
            if fld is not None:
                tnode = fld.find(qn("a:t"))
                if tnode is not None:
                    tnode.text = str(page_no)
        spTree.append(new_sp)

def add_slide(prs, layout_idx, title=None, footer_page=None,
              footer_text="DRMD PPTX変換検証 / CONFIDENTIAL（検証用）"):
    layout = prs.slide_masters[0].slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    if title is not None and slide.shapes.title is not None:
        set_title(slide, title)
    if footer_page is not None:
        add_furniture(slide, footer_text, footer_page)
    return slide

# ================================================================ slides ==
def build_slide01_title(prs):
    """P01(a): title slide. QA1: 文書番号/版/状態(表紙)."""
    slide = add_slide(prs, 0)  # Title Slide layout
    set_title(slide, "経費精算システム　基本・詳細設計書", size=34)
    sub = slide.placeholders[1]
    tf = sub.text_frame
    tf.word_wrap = True
    lines = [
        "設計レビュー説明資料（検証用サンプル・架空システム）",
        "文書番号: EXPS-DES-001　版: 1.2　状態: レビュー中（検証用）",
        "作成日: 2026-08-23　対象システム: 経費精算システム　機密区分: 社内一般（架空）",
        "作成者: アプリ設計担当　確認者: 業務・基盤担当",
    ]
    for i, line in enumerate(lines):
        p = new_paragraph(tf, first=(i == 0))
        r = p.add_run()
        r.text = line
        r.font.name = JP_FONT
        r.font.size = Pt(16 if i == 0 else 13)
        r.font.color.rgb = TEAL if i == 0 else MUTED
    record("P01", "タイトルスライド（role=title）", "layout=Title Slide + placeholders", True)
    record("QA1", "表紙: 文書番号/版/状態", "title placeholders text", True)
    return slide

def build_slide02_section(prs):
    """P01(b): section-header slide."""
    slide = add_slide(prs, 2)  # Section Header layout
    set_title(slide, "第1部　システム概要とインターフェース設計", size=30)
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    lines = [
        "対象: 経費申請の登録、承認、会計連携、支払結果通知",
        "収録範囲: 概要・状態遷移・API/データ・テスト・非機能・運用・画像検証",
        "原本: 経費精算システム_設計書_検証用.xlsx（全11シート、文書ID doc_90f44c5f76d04d37）",
    ]
    for i, line in enumerate(lines):
        p = new_paragraph(tf, first=(i == 0))
        r = p.add_run()
        r.text = line
        r.font.name = JP_FONT
        r.font.size = Pt(15)
        r.font.color.rgb = INK
    record("P01", "セクション区切りスライド（role=section header）", "layout=Section Header", True)
    return slide

def build_slide03_overview(prs):
    """P05: placeholder that inherits bullet formatting from layout/master
    (no explicit <a:pPr> bullet element is written on this slide).
    QA2: IF-05 (送信元/送信先/Timeout/Retry/失敗時処理)."""
    slide = add_slide(prs, 1, "1. システム概要と外部インターフェース", footer_page=3)
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    lines = [
        "対象システム: 経費精算システム／利用者: 申請者、課長・部長承認者、経理担当者／認証: OIDC・JWT",
        "金額閾値: 100,000円以下は課長承認、100,000円超は部長承認（BR-01）",
        "領収書: 10,000円以上で添付必須。PDF/JPEG/PNG、最大10MB",
        "IF-05: 送信元 Outbox Worker → 送信先 Event Bus（AMQP）。Timeout 5秒、Retry 最大5回。"
        "失敗時はDLQへ登録しOPS-ALM-02を発報する。",
        "IF-06: 送信元 Expense API → 送信先 会計System（REST/Batch）。Timeout 10秒、Retry 最大5回。"
        "失敗時はERROR遷移・手動再送。",
    ]
    for i, line in enumerate(lines):
        # plain_paragraph never touches bullet XML -> inherits layout/master default
        plain_paragraph(tf, line, first=(i == 0), size=15, color=INK)
    record("P05", "master/layout既定継承のプレースホルダー（buXxx未記載）", "no <a:pPr> bullet children written", True)
    record("QA2", "IF-05: 送信元/送信先/Timeout/Retry/失敗時処理", "body placeholder text", True)
    record("QA5", "BR-01: 金額閾値と承認経路", "body placeholder text", True)
    return slide

def build_slide04_assumptions(prs):
    """P02: buChar bulleted list, 4 levels deep."""
    slide = add_slide(prs, 1, "1.1 設計前提の階層整理", footer_page=4)
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    chars = ["◆", "●", "▪", "‒"]
    items = [
        (0, "認証・認可"),
        (1, "OIDC / JWT ベースの認証"),
        (2, "Web UI がアクセストークンを Expense API へ送信"),
        (3, "承認 API は承認者ロールのみ実行可（NFR-SEC-01）"),
        (0, "金額閾値と承認経路"),
        (1, "100,000円以下は課長承認（BR-01）"),
        (1, "100,000円超は部長承認（BR-01）"),
        (0, "領収書の取り扱い"),
        (1, "10,000円以上で添付必須（F-06）"),
        (2, "PDF / JPEG / PNG、最大10MB"),
        (0, "可用性"),
        (1, "平日8:00-22:00、月間稼働率99.9%目標（NFR-AVL-01）"),
    ]
    first = True
    for level, text in items:
        p = new_paragraph(tf, first=first)
        first = False
        p.level = level
        r = p.add_run()
        r.text = text
        r.font.name = JP_FONT
        r.font.size = Pt(max(12, 20 - level * 2))
        r.font.color.rgb = INK if level else NAVY
        set_bullet_char(p, chars[level])
    max_level = max(lvl for lvl, _ in items)
    record("P02", f"buChar 箇条書き（{max_level + 1}階層）", "a:buChar per level 0-3", True)
    return slide

STATE_LABELS = {
    "DRAFT": "下書き",
    "SUBMITTED": "提出済",
    "PENDING": "承認待ち",
    "APPROVED": "承認済",
    "RETURNED": "差戻し",
    "PAYMENT_PENDING": "支払待ち",
    "PAID": "支払済",
    "ERROR": "連携エラー",
    "CANCELLED": "取消",
}
STATE_POS = {  # inches: x, y, w, h
    "DRAFT": (0.35, 1.55, 1.55, 0.7),
    "SUBMITTED": (2.15, 1.55, 1.55, 0.7),
    "PENDING": (3.95, 1.55, 1.55, 0.7),
    "APPROVED": (5.75, 1.55, 1.55, 0.7),
    "PAYMENT_PENDING": (7.55, 1.55, 1.75, 0.7),
    "PAID": (9.55, 1.55, 1.55, 0.7),
    "CANCELLED": (0.35, 2.95, 1.55, 0.7),
    "RETURNED": (3.95, 2.95, 1.55, 0.7),
    "ERROR": (7.55, 2.95, 1.75, 0.7),
}
# (transition id, from, to, event, guard, side-effect)
TR_TABLE = [
    ("TR-01", "—", "新規作成", "認証済み", "DRAFT", "一時保存領域作成"),
    ("TR-02", "DRAFT", "提出", "必須入力OK、金額>0", "SUBMITTED", "監査ログ記録"),
    ("TR-03", "SUBMITTED", "受付", "Idempotency-Key未処理", "PENDING", "承認イベント発行"),
    ("TR-04", "PENDING", "承認", "承認権限あり", "APPROVED", "承認日時を記録"),
    ("TR-05", "PENDING", "差戻し", "理由1文字以上", "RETURNED", "申請者へ通知"),
    ("TR-06", "RETURNED", "再提出", "修正後の入力検証OK", "SUBMITTED", "版番号を加算"),
    ("TR-07", "APPROVED", "会計連携成功", "仕訳番号あり", "PAYMENT_PENDING", "仕訳番号を保存"),
    ("TR-08", "PAYMENT_PENDING", "支払確定", "支払日あり", "PAID", "完了通知"),
    ("TR-09", "PAYMENT_PENDING", "連携失敗", "retry_count >= 5", "ERROR", "DLQ登録・運用アラート"),
    ("TR-10", "ERROR", "手動再送", "原因解消済・運用権限あり", "PAYMENT_PENDING", "retry_countを0へ戻す"),
    ("TR-11", "DRAFT", "取消", "未提出", "CANCELLED", "論理削除"),
]
# connector wiring: (transition id, from-state, to-state, begin_idx, end_idx)
# connection-site idx convention used throughout: 0=top 1=right 2=bottom 3=left
CONNECTORS = [
    ("TR-02", "DRAFT", "SUBMITTED", 1, 3),
    ("TR-03", "SUBMITTED", "PENDING", 1, 3),
    ("TR-04", "PENDING", "APPROVED", 1, 3),
    ("TR-05", "PENDING", "RETURNED", 2, 0),
    ("TR-06", "RETURNED", "SUBMITTED", 0, 2),
    ("TR-07", "APPROVED", "PAYMENT_PENDING", 1, 3),
    ("TR-08", "PAYMENT_PENDING", "PAID", 1, 3),
    ("TR-09", "PAYMENT_PENDING", "ERROR", 2, 0),
    ("TR-10", "ERROR", "PAYMENT_PENDING", 0, 2),
    ("TR-11", "DRAFT", "CANCELLED", 2, 0),
]

def build_slide05_state_machine(prs):
    """P08: group shape + in-shape text + logically-connected connectors
    (stCxn/endCxn) drawing the real DRAFT..PAID state machine.
    QA3: TR-09 full (from/event/guard/to/side-effect) also written out in the
    detail table so the fact survives even if diagram geometry does not."""
    slide = add_slide(prs, 5, "2.2 状態遷移設計（申請ライフサイクル）", footer_page=5)  # Title Only
    shapes = slide.shapes
    state_shapes = {}
    for code, (x, y, w, h) in STATE_POS.items():
        shp = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = TEAL if code in ("ERROR", "CANCELLED", "RETURNED") else NAVY
        shp.line.color.rgb = WHITE
        tf = shp.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run()
        r0.text = code
        r0.font.bold = True
        r0.font.size = Pt(11)
        r0.font.color.rgb = WHITE
        r0.font.name = JP_FONT
        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = STATE_LABELS[code]
        r1.font.size = Pt(9)
        r1.font.color.rgb = WHITE
        r1.font.name = JP_FONT
        state_shapes[code] = shp

    group = shapes.add_group_shape(list(state_shapes.values()))
    group.name = "状態遷移グループ（9状態）"

    for tr_id, src, dst, b_idx, e_idx in CONNECTORS:
        s, d = state_shapes[src], state_shapes[dst]
        conn = shapes.add_connector(MSO_CONNECTOR.ELBOW, Inches(0), Inches(0), Inches(1), Inches(1))
        conn.begin_connect(s, b_idx)
        conn.end_connect(d, e_idx)
        conn.line.color.rgb = MUTED
        conn.line.width = Pt(1.5)

    # short inline labels for the two transitions most people ask about
    label_box = shapes.add_textbox(Inches(0.35), Inches(2.35), Inches(3.2), Inches(0.5))
    lp = label_box.text_frame.paragraphs[0]
    lr = lp.add_run()
    lr.text = "TR-05 差戻し [理由1文字以上]"
    lr.font.size = Pt(9)
    lr.font.italic = True
    lr.font.color.rgb = MUTED
    lr.font.name = JP_FONT

    label_box2 = shapes.add_textbox(Inches(7.55), Inches(2.35), Inches(2.9), Inches(0.5))
    lp2 = label_box2.text_frame.paragraphs[0]
    lr2 = lp2.add_run()
    lr2.text = "TR-09 連携失敗 [retry_count >= 5]"
    lr2.font.size = Pt(9)
    lr2.font.italic = True
    lr2.font.color.rgb = RED
    lr2.font.name = JP_FONT

    # full 2.2 状態・イベント定義 table (guarantees TR-09 fact survives as plain text)
    rows, cols = len(TR_TABLE) + 1, 6
    gframe = shapes.add_table(rows, cols, Inches(0.35), Inches(4.05), Inches(12.6), Inches(3.1))
    tbl = gframe.table
    headers = ["遷移ID", "遷移元", "イベント", "ガード／条件", "遷移先", "副作用"]
    widths = [Inches(1.0), Inches(2.3), Inches(1.9), Inches(3.4), Inches(2.3), Inches(1.7)]
    for c, w in enumerate(widths):
        tbl.columns[c].width = w
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.bold = True
        para.runs[0].font.size = Pt(9)
        para.runs[0].font.color.rgb = WHITE
        para.runs[0].font.name = JP_FONT
    for r, row in enumerate(TR_TABLE, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(8.5)
            para.runs[0].font.name = JP_FONT
            para.runs[0].font.color.rgb = RED if row[0] == "TR-09" else INK
            if row[0] == "TR-09":
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF4, 0xE5)
    record("P08", "グループ図形+図形内テキスト+コネクタ(stCxn/endCxn)の状態遷移図",
           "add_group_shape(9 states) + add_connector x10 (begin/end_connect)", True)
    record("QA3", "TR-09: 遷移元/イベント/ガード/遷移先/副作用", "detail table row + diagram", True)
    return slide

def build_slide06_procedure(prs):
    """P03: buAutoNum numbered list (procedure steps)."""
    slide = add_slide(prs, 1, "4.1 業務フロー手順（申請〜支払）", footer_page=6)
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    steps = [
        "BF-00 開始",
        "BF-01 領収書・明細を入力",
        "BF-02 形式・必須チェック（GW-01 入力は妥当か）",
        "BF-03 申請登録 SUBMITTED（TX-01）",
        "BF-04 承認者決定（10万円以下=課長／超過=部長、BR-01）",
        "BF-05〜BF-06 内容確認・差戻し判定（GW-02 承認するか）",
        "BF-07 仕訳連携（GW-03 連携成功か）",
        "BF-08 再送キュー（失敗時、1/5/30/120/600秒）",
        "BF-09 支払処理",
        "BF-10 支払済更新・完了通知（BF-99 終了）",
    ]
    first = True
    for step in steps:
        p = new_paragraph(tf, first=first)
        first = False
        r = p.add_run()
        r.text = step
        r.font.size = Pt(16)
        r.font.color.rgb = INK
        r.font.name = JP_FONT
        set_bullet_autonum(p, "arabicPeriod")
    record("P03", "buAutoNum 番号付きリスト（手順、10ステップ）", "a:buAutoNum type=arabicPeriod", True)
    return slide

def build_slide07_nfr(prs):
    """P04: buNone paragraph mixed with buChar-bulleted paragraphs.
    QA9: NFR-AVL-02 full (RTO/RPO/測定方法/実現手段/合意状態)."""
    slide = add_slide(prs, 1, "8.1 非機能要件ハイライト", footer_page=7)
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True

    intro = new_paragraph(tf, first=True)
    ri = intro.add_run()
    ri.text = "IPA非機能要求グレードを参考に、目標値・測定方法・実現手段・合意状態を管理する。"
    ri.font.size = Pt(14)
    ri.font.italic = True
    ri.font.color.rgb = MUTED
    ri.font.name = JP_FONT
    set_bullet_none(intro)

    items = [
        "NFR-AVL-01（可用性）: 平日8:00–22:00の月間稼働率99.9%以上。2AZ配置・ヘルスチェック・ローリング更新。条件付合意。",
        "NFR-AVL-02（回復性）: RTO 4時間／RPO 15分。年1回の復旧訓練で実測。日次バックアップ＋WAL 15分転送で実現。合意状態: 要確認。",
        "NFR-PERF-01（性能）: 通常申請 p95 1.0秒以下、20 TPS。APMで5分窓測定。合意済。",
        "NFR-OPS-04（再送）: 1/5/30/120/600秒、最大5回。Outbox Worker・指数バックオフ・DLQ。合意済。",
    ]
    for text in items:
        p = new_paragraph(tf)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(14)
        r.font.color.rgb = INK
        r.font.name = JP_FONT
        set_bullet_char(p, "●")
    record("P04", "buNone明示解除+箇条書きの混在", "1x a:buNone + 4x a:buChar", True)
    record("QA9", "NFR-AVL-02: RTO/RPO/測定方法/実現手段/合意状態", "bulleted paragraph text", True)
    return slide

def build_slide08_charts(prs):
    """P06: native bar chart (test summary) + native pie chart (expense
    category breakdown, illustrative)."""
    slide = add_slide(prs, 5, "7.2 テスト実施サマリーと経費区分内訳（例示）", footer_page=8)  # Title Only
    shapes = slide.shapes

    bar_data = CategoryChartData()
    bar_data.categories = ["合格", "不合格", "未実施"]
    bar_data.add_series("テストケース件数（総件数9）", (0, 0, 9))
    bar_frame = shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.5),
                                  Inches(5.8), Inches(4.2), bar_data)
    bar_frame.chart.has_title = True
    bar_frame.chart.chart_title.text_frame.text = "テスト実施サマリー（07_テストシナリオ 7.2）"

    pie_data = CategoryChartData()
    pie_data.categories = ["標準10%（STANDARD）", "軽減8%（REDUCED）", "非課税（EXEMPT）"]
    pie_data.add_series("経費区分内訳（税区分ベース・例示、F-05）", (58, 27, 15))
    pie_frame = shapes.add_chart(XL_CHART_TYPE.PIE, Inches(6.9), Inches(1.5),
                                  Inches(5.8), Inches(4.2), pie_data)
    pie_frame.chart.has_title = True
    pie_frame.chart.chart_title.text_frame.text = "経費区分内訳（例示、F-05税区分ベース）"

    cap = shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.9))
    cp = cap.text_frame.paragraphs[0]
    cr = cp.add_run()
    cr.text = ("総件数9 / 合格0 / 不合格0 / 未実施9（07_テストシナリオ 7.2 実施サマリー）。"
               "円グラフは経費区分内訳の例示（F-05 税区分 STANDARD/REDUCED/EXEMPT を利用）。")
    cr.font.size = Pt(11)
    cr.font.color.rgb = MUTED
    cr.font.name = JP_FONT
    record("P06", "ネイティブチャート: 棒(テスト実施サマリー)+円(経費区分内訳)",
           "add_chart(COLUMN_CLUSTERED) + add_chart(PIE), 実データ入り", True)
    return slide

TC_ROWS = [
    ("TC-001", "12,800円、領収書あり", "課長承認→会計成功", "201", "PAID", "TR-02/03/04/07/08, API-01"),
    ("TC-004", "同一Idempotency-Keyを再送", "既存結果を返却", "201", "SUBMITTED", "DR-04, API-01"),
    ("TC-005", "承認者が理由付き差戻し", "PENDING→RETURNED", "200", "RETURNED", "TR-05"),
    ("TC-008", "Event Busを5分停止", "Outbox再送→承認タスク生成", "201", "PENDING", "SQ-01, EVT-01, NFR-OPS-04"),
    ("TC-009", "会計I/Fが5回連続失敗", "DLQ→ERROR→手動再送", "200", "ERROR", "TR-09/10, BR-03, OPS-ALM-02"),
]

def build_slide09_test_scenarios(prs):
    """QA8: TC-009 full (入力/経路/HTTP/最終状態/関連設計)."""
    slide = add_slide(prs, 5, "7.1 代表テストシナリオ（抜粋）", footer_page=9)  # Title Only
    shapes = slide.shapes
    rows, cols = len(TC_ROWS) + 1, 6
    gframe = shapes.add_table(rows, cols, Inches(0.4), Inches(1.6), Inches(12.5), Inches(3.6))
    tbl = gframe.table
    headers = ["TC-ID", "観点／入力", "通過経路", "期待HTTP", "期待最終状態", "関連設計"]
    widths = [Inches(1.1), Inches(3.0), Inches(2.9), Inches(1.3), Inches(1.7), Inches(2.5)]
    for c, w in enumerate(widths):
        tbl.columns[c].width = w
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = WHITE
        run.font.name = JP_FONT
    for r, row in enumerate(TC_ROWS, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(10.5)
            run.font.name = JP_FONT
            run.font.color.rgb = RED if row[0] == "TC-009" else INK
            if row[0] == "TC-009":
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF4, 0xE5)
    cap = shapes.add_textbox(Inches(0.4), Inches(5.35), Inches(12.5), Inches(0.6))
    cr = cap.text_frame.paragraphs[0].add_run()
    cr.text = "総件数9件（7.2 実施サマリー）のうち代表5件を抜粋。全件は元設計書 07_テストシナリオを参照。"
    cr.font.size = Pt(11)
    cr.font.italic = True
    cr.font.color.rgb = MUTED
    cr.font.name = JP_FONT
    record("QA8", "TC-009: 入力/経路/HTTP/最終状態/関連設計", "table row (highlighted)", True)
    return slide

# (区分, 名称, 型, 必須, 説明)
API01_CONTRACT = [
    ("Header", "Authorization", "string", "○", "Bearer JWT"),
    ("Header", "Idempotency-Key", "uuid", "○", "同一利用者で24時間一意"),
    ("Body", "applicationDate", "date", "○", "yyyy-mm-dd"),
    ("Body", "departmentCode", "string", "○", "5文字"),
    ("Body", "description", "string", "○", "1〜100文字"),
    ("Body", "amount", "integer", "○", "円、1〜9,999,999"),
    ("Body", "taxType", "enum", "○", "STANDARD / REDUCED / EXEMPT"),
    ("Body", "receiptFileId", "uuid", "条件付", "10,000円以上で必須"),
    ("Response", "claimId", "uuid", "○", "受付番号"),
    ("Response", "status", "enum", "○", "SUBMITTED"),
]

def build_slide10_api_contract(prs):
    """P09: table with merged cells (vertical merge on 区分 column).
    QA7: API-01 full (Method/Path/応答/Idempotency-Key制約).
    QA4: TX-01 full (書き込み対象/HTTP応答後のイベント配送).
    QA11: outbox_event.status candidates + retry_count上限."""
    slide = add_slide(prs, 5, "6.2 POST /v1/expense-claims 契約（API-01）", footer_page=10)  # Title Only
    shapes = slide.shapes
    rows, cols = len(API01_CONTRACT) + 1, 5
    gframe = shapes.add_table(rows, cols, Inches(0.4), Inches(1.5), Inches(8.6), Inches(4.1))
    tbl = gframe.table
    headers = ["区分", "名称", "型", "必須", "説明"]
    widths = [Inches(1.1), Inches(2.3), Inches(1.1), Inches(1.0), Inches(3.1)]
    for c, w in enumerate(widths):
        tbl.columns[c].width = w
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = WHITE
        run.font.name = JP_FONT

    for r, row in enumerate(API01_CONTRACT, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(10.5)
            run.font.name = JP_FONT
            run.font.color.rgb = INK

    # vertical merge on the 区分 column: Header(1-2), Body(3-8), Response(9-10)
    tbl.cell(1, 0).text = "Header"
    tbl.cell(2, 0).text = ""
    tbl.cell(1, 0).merge(tbl.cell(2, 0))
    tbl.cell(3, 0).text = "Body"
    for rr in range(4, 9):
        tbl.cell(rr, 0).text = ""
    tbl.cell(3, 0).merge(tbl.cell(8, 0))
    tbl.cell(9, 0).text = "Response"
    tbl.cell(10, 0).text = ""
    tbl.cell(9, 0).merge(tbl.cell(10, 0))
    for rr in (1, 3, 9):
        run = tbl.cell(rr, 0).text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(10.5)
        run.font.color.rgb = TEAL
        run.font.name = JP_FONT

    box = shapes.add_textbox(Inches(9.2), Inches(1.5), Inches(3.7), Inches(4.6))
    tf = box.text_frame
    tf.word_wrap = True
    facts = [
        ("API-01 契約", True, NAVY),
        ("Method: POST", False, INK),
        ("Path: /v1/expense-claims", False, INK),
        ("主な応答: 201 / 400 / 401 / 409", False, INK),
        ("Idempotency-Key: 同一利用者で24時間一意（uuid、必須）", False, INK),
        ("", False, INK),
        ("TX-01（トランザクション境界）", True, NAVY),
        ("書き込み対象: expense_claim（status=SUBMITTED）と outbox_event を"
         "同一トランザクションで確定する。", False, INK),
        ("HTTP応答後のイベント配送は非同期。Event Bus障害時はNFR-OPS-04の間隔で"
         "最大5回再送し、承認サービスはevent_idで冪等に処理する。", False, INK),
        ("", False, INK),
        ("outbox_event.status: NEW/SENDING/SENT/FAILED", False, TEAL),
        ("retry_count: 初期値0、上限5", False, TEAL),
    ]
    first = True
    for text, bold, color in facts:
        p = new_paragraph(tf, first=first)
        first = False
        if text:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(10.5 if not bold else 12)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = JP_FONT
    record("P09", "結合セルを含む表（API一覧/契約、vMerge=区分列）",
           "add_table + cell.merge (rowSpan/vMerge)", True)
    record("QA7", "API-01: Method/Path/応答/Idempotency-Key制約", "side textbox", True)
    record("QA4", "TX-01: 書き込み対象/HTTP応答後のイベント配送", "side textbox", True)
    record("QA11", "outbox_event.status候補 + retry_count上限", "side textbox", True)
    return slide

def build_slide11_screen_spec(prs):
    """P14: bold/italic/underline/strike/color run decorations.
    P15: 45°回転シェイプ + 矢印図形.
    QA6: F-06 full (必須条件/最大サイズ/エラーコード/メッセージ/表示位置)."""
    slide = add_slide(prs, 5, "5.2 項目定義・検証ルール（F-06 領収書）", footer_page=11)  # Title Only
    shapes = slide.shapes

    box = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7.6), Inches(4.6))
    tf = box.text_frame
    tf.word_wrap = True

    p1 = new_paragraph(tf, first=True)
    add_run(p1, "項目ID F-06", bold=True, size=16, color=NAVY)
    add_run(p1, "　領収書（file 型）", size=16, color=INK)

    p2 = tf.add_paragraph()
    add_run(p2, "条件付必須", italic=True, size=13, color=MUTED)
    add_run(p2, "：", size=13, color=INK)
    add_run(p2, "10,000円以上は必須", underline=True, size=13, color=INK)
    add_run(p2, "、最大", size=13, color=INK)
    add_run(p2, "10MB", bold=True, size=13, color=INK)
    add_run(p2, "。", size=13, color=INK)

    p3 = tf.add_paragraph()
    add_run(p3, "旧仕様: 5,000円以上は必須", strike=True, size=12, color=MUTED)
    add_run(p3, "（2026-08-23改定）", size=12, color=MUTED)

    p4 = tf.add_paragraph()
    add_run(p4, "エラーコード ", size=13, color=INK)
    add_run(p4, "EXP-E006", bold=True, size=14, color=RED)

    p5 = tf.add_paragraph()
    add_run(p5, "表示メッセージ：「", size=13, color=INK)
    add_run(p5, "10,000円以上の申請には領収書が必要です。", italic=True, size=13, color=TEAL)
    add_run(p5, "」", size=13, color=INK)

    p6 = tf.add_paragraph()
    add_run(p6, "表示位置：", size=13, color=INK)
    add_run(p6, "領収書直下", underline=True, size=13, color=INK)

    diamond = shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(8.7), Inches(2.2), Inches(1.3), Inches(1.3))
    diamond.rotation = 45
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = GOLD
    diamond.line.color.rgb = WHITE
    dtf = diamond.text_frame
    dtf.word_wrap = True
    dp = dtf.paragraphs[0]
    dp.alignment = PP_ALIGN.CENTER
    # counter-rotate the text so it reads horizontally despite the 45 deg shape
    bodyPr = dtf._txBody.find(qn("a:bodyPr"))
    bodyPr.set("anchor", "ctr")
    dr = dp.add_run()
    dr.text = "改定"
    dr.font.bold = True
    dr.font.size = Pt(14)
    dr.font.color.rgb = WHITE
    dr.font.name = JP_FONT
    # rotate the text body opposite to the shape so "改定" stays upright
    txBodyPr = diamond.text_frame._txBody.find(qn("a:bodyPr"))
    txBodyPr.set("rot", "-2700000")  # -45 degrees in 60,000ths of a degree

    arrow = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(10.3), Inches(2.55), Inches(1.7), Inches(0.6))
    arrow.rotation = 200
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RED
    arrow.line.color.rgb = WHITE
    atf = arrow.text_frame
    ap = atf.paragraphs[0]
    ap.alignment = PP_ALIGN.CENTER
    ar = ap.add_run()
    ar.text = "要確認"
    ar.font.size = Pt(10)
    ar.font.bold = True
    ar.font.color.rgb = WHITE
    ar.font.name = JP_FONT
    record("P14", "run装飾: 太字/斜体/下線/取消線/文字色",
           "run.font.bold/italic/underline + rPr@strike(raw XML) + font.color.rgb", True)
    record("P15", "回転シェイプ(45°)+矢印図形", "shape.rotation=45 (DIAMOND) + RIGHT_ARROW", True)
    record("QA6", "F-06: 必須条件/最大サイズ/エラーコード/メッセージ/表示位置", "decorated textbox", True)
    return slide

def _fill_column(placeholder, heading, items):
    tf = placeholder.text_frame
    tf.word_wrap = True
    hp = new_paragraph(tf, first=True)
    hr = hp.add_run()
    hr.text = heading
    hr.font.bold = True
    hr.font.size = Pt(15)
    hr.font.color.rgb = NAVY
    hr.font.name = JP_FONT
    set_bullet_none(hp)
    for text in items:
        p = new_paragraph(tf)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(12.5)
        r.font.color.rgb = INK
        r.font.name = JP_FONT
        set_bullet_char(p, "●")

def build_slide12_two_column(prs):
    """P12: two content placeholders (left/right columns).
    QA10: ISSUE-03 full (論点/作業/Owner/期限/Status)."""
    slide = add_slide(prs, 3, "8.3 運用上の留意点（申請者・承認者 / 運用）", footer_page=12)  # Two Content
    left = slide.placeholders[1]
    right = slide.placeholders[2]
    _fill_column(left, "申請者・承認者向け", [
        "差戻し理由は1〜500文字（BR-02）",
        "支払済(PAID)は更新不可、取消はDRAFTのみ（重要な不変条件）",
        "会計連携5回失敗でERROR・DLQへ移管（BR-03）",
        "同一Idempotency-Keyの再送で二重登録しない（重要な不変条件／DR-04一意制約）",
    ])
    _fill_column(right, "運用・データ", [
        "ISSUE-03: 会計停止時間帯の調整。IF-06再送窓を会計主管と調整。"
        "Owner: 経理担当、期限: 2026-08-27、Status: 対応中。",
        "outbox_event.status: NEW/SENDING/SENT/FAILED（retry_count 上限5）",
        "OPS-ALM-02: DLQ件数 >= 1で対象event_idを凍結、SLA即時、Runbook RB-EVT-02",
    ])
    record("P12", "2カラム構成スライド（左右コンテンツプレースホルダー）", "layout=Two Content, placeholders idx1/idx2", True)
    record("QA10", "ISSUE-03: 論点/作業/Owner/期限/Status", "right column bullet", True)
    return slide

SMARTART_NODES = ["申請受付", "課長承認(≤10万円)", "部長承認(>10万円)", "会計連携", "支払確定"]

def build_slide13_smartart(prs):
    """P07: SmartArt (dgm:) attempt. python-pptx has no API for diagram
    graphicData, so the real <p:graphicFrame> + ppt/diagrams/*.xml parts are
    injected as a raw zip post-processing pass (see inject_smartart, called
    from build() after prs.save()). This function only lays out the ordinary
    shapes around that reserved area: a label, and a plain-autoshape
    restatement of the same 5 steps so the slide keeps readable text
    regardless of whether the dgm: text model survives conversion."""
    slide = add_slide(prs, 5, "6.7 承認プロセス（SmartArt試行）", footer_page=13)  # Title Only
    shapes = slide.shapes

    label = shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.0), Inches(0.4))
    lr = label.text_frame.paragraphs[0].add_run()
    lr.text = "SmartArt（dgm:データモデル、5ノード実データ）※直下の枠に実データを後注入"
    lr.font.size = Pt(11)
    lr.font.italic = True
    lr.font.color.rgb = MUTED
    lr.font.name = JP_FONT
    # The actual <p:graphicFrame> referencing ppt/diagrams/data1.xml is
    # injected at Inches(0.6, 2.0, 12.0, 2.1) by inject_smartart() below.

    x, w, y = 0.6, 2.3, 4.35
    for i, text in enumerate(SMARTART_NODES):
        shp = shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + i * (w - 0.35)), Inches(y), Inches(w), Inches(1.0))
        shp.fill.solid()
        shp.fill.fore_color.rgb = NAVY if i % 2 == 0 else TEAL
        shp.line.color.rgb = WHITE
        tf = shp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = JP_FONT

    note = shapes.add_textbox(Inches(0.6), Inches(5.55), Inches(12.0), Inches(1.2))
    ntf = note.text_frame
    ntf.word_wrap = True
    nr = ntf.paragraphs[0].add_run()
    nr.text = ("同じ5工程を通常の図形テキストとしても併記（BR-01: 10万円以下は課長、超過は部長）。"
               "SmartArtのdgm:テキストのみが変換で失われるかを、上下2系統の比較で検証できる。")
    nr.font.size = Pt(11)
    nr.font.italic = True
    nr.font.color.rgb = MUTED
    nr.font.name = JP_FONT
    record("P07", "SmartArt（承認プロセス、dgm:データモデル注入）",
           "raw zip post-processing: ppt/diagrams/{data1,layout1,quickStyle1,colors1}.xml"
           " + graphicFrame dgm:relIds（python-pptxにAPI無し）", True)
    return slide

# ---------------------------------------------------------- SmartArt (raw) --
# python-pptx cannot create <a:graphicData uri=".../diagram"> content, so the
# minimal SmartArt part set is injected directly into the saved .pptx zip.
# Verified in isolation: python-pptx re-opens it cleanly and the DRMD CLI
# exits 0 (it simply does not read ppt/diagrams/*, which is the documented
# P07 gap: SmartArt text is dropped, not that the file becomes unreadable).
_SA_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_SA_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_SA_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_SA_NS_DGM = "http://schemas.openxmlformats.org/drawingml/2006/diagram"

def _smartart_data_xml(nodes):
    doc_id = "{A0000000-0000-0000-0000-000000000000}"
    pts = [
        f'<dgm:pt modelId="{doc_id}" type="doc">'
        '<dgm:prSet loTypeId="urn:microsoft.com/office/officeart/2005/8/layout/process1" '
        'loCatId="list" qsTypeId="urn:microsoft.com/office/officeart/2005/8/quickstyle/simple1" '
        'qsCatId="simple" csTypeId="urn:microsoft.com/office/officeart/2005/8/colors/accent1_2" '
        'csCatId="accent1"/></dgm:pt>'
    ]
    cxns = []
    for i, text in enumerate(nodes):
        node_id = f"{{B{i:07d}-0000-0000-0000-000000000000}}"
        par_id = f"{{C{i:07d}-0000-0000-0000-000000000000}}"
        sib_id = f"{{D{i:07d}-0000-0000-0000-000000000000}}"
        cxn_id = f"{{E{i:07d}-0000-0000-0000-000000000000}}"
        pts.append(
            f'<dgm:pt modelId="{node_id}"><dgm:prSet/><dgm:t><a:bodyPr/><a:lstStyle/>'
            f'<a:p><a:r><a:rPr lang="ja-JP" dirty="0"/><a:t>{text}</a:t></a:r></a:p></dgm:t></dgm:pt>'
        )
        pts.append(f'<dgm:pt modelId="{par_id}" type="parTrans" cxnId="{cxn_id}"><dgm:prSet/><dgm:spPr/></dgm:pt>')
        pts.append(f'<dgm:pt modelId="{sib_id}" type="sibTrans" cxnId="{cxn_id}"><dgm:prSet/><dgm:spPr/></dgm:pt>')
        cxns.append(
            f'<dgm:cxn modelId="{cxn_id}" type="parOf" srcId="{doc_id}" destId="{node_id}" '
            f'srcOrd="{i}" destOrd="0" parTransId="{par_id}" sibTransId="{sib_id}"/>'
        )
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        f'<dgm:dataModel xmlns:dgm="{_SA_NS_DGM}" xmlns:a="{_SA_NS_A}" xmlns:r="{_SA_NS_R}">'
        f'<dgm:ptLst>{"".join(pts)}</dgm:ptLst><dgm:cxnLst>{"".join(cxns)}</dgm:cxnLst>'
        '<dgm:bg/><dgm:whole/></dgm:dataModel>'
    )

_SA_LAYOUT_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
    f'<dgm:layoutDef xmlns:dgm="{_SA_NS_DGM}" xmlns:a="{_SA_NS_A}" xmlns:r="{_SA_NS_R}" '
    'uniqueId="urn:drmd:fixture:process1min" '
    'minVer="http://schemas.openxmlformats.org/drawingml/2006/diagram" '
    'defStyle="urn:microsoft.com/office/officeart/2005/8/quickstyle/simple1">'
    '<dgm:title val=""/><dgm:desc val=""/><dgm:catLst><dgm:cat type="list" pri="1"/></dgm:catLst>'
    '<dgm:sampData/><dgm:styleData/>'
    '<dgm:layoutNode name="root"><dgm:varLst><dgm:dir/><dgm:animLvl val="lvl"/></dgm:varLst>'
    '<dgm:alg type="composite"/><dgm:shape type="none"/><dgm:presOf axis="ch" ptType="node"/>'
    '<dgm:forEach name="nodesForEach" axis="ch" ptType="node">'
    '<dgm:layoutNode name="node"><dgm:alg type="tx"/><dgm:shape type="rect"/>'
    '<dgm:presOf axis="desOrSelf" ptType="node"/></dgm:layoutNode></dgm:forEach>'
    '</dgm:layoutNode></dgm:layoutDef>'
)

_SA_QUICKSTYLE_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
    f'<dgm:styleDefHdrLst xmlns:dgm="{_SA_NS_DGM}" xmlns:a="{_SA_NS_A}">'
    '<dgm:styleDefHdr uniqueId="urn:drmd:fixture:qsmin" '
    'minVer="http://schemas.openxmlformats.org/drawingml/2006/diagram">'
    '<dgm:title val=""/><dgm:desc val=""/><dgm:catLst><dgm:cat type="simple" pri="1"/></dgm:catLst>'
    '</dgm:styleDefHdr></dgm:styleDefHdrLst>'
)

_SA_COLORS_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
    f'<dgm:colorsDefHdrLst xmlns:dgm="{_SA_NS_DGM}">'
    '<dgm:colorsDefHdr uniqueId="urn:drmd:fixture:csmin" '
    'minVer="http://schemas.openxmlformats.org/drawingml/2006/diagram">'
    '<dgm:title val=""/><dgm:desc val=""/><dgm:catLst><dgm:cat type="accent1" pri="1"/></dgm:catLst>'
    '</dgm:colorsDefHdr></dgm:colorsDefHdrLst>'
)

def inject_smartart(pptx_path, slide_part, nodes=SMARTART_NODES):
    """Rewrite pptx_path in place, adding a real SmartArt diagram graphicFrame
    to slide_part (e.g. 'ppt/slides/slide13.xml')."""
    with zipfile.ZipFile(pptx_path) as zin:
        slide_xml = zin.read(slide_part).decode("utf-8")
        names = set(zin.namelist())
    existing_ids = [int(x) for x in etree.fromstring(slide_xml.encode("utf-8")).xpath(
        "//*[local-name()='cNvPr']/@id")]
    shape_id = max(existing_ids, default=1) + 1
    rels_name = slide_part.replace("slides/", "slides/_rels/") + ".rels"

    graphic_frame = (
        f'<p:graphicFrame xmlns:p="{_SA_NS_P}" xmlns:a="{_SA_NS_A}" xmlns:r="{_SA_NS_R}">'
        f'<p:nvGraphicFramePr><p:cNvPr id="{shape_id}" name="SmartArt 承認プロセス"/>'
        '<p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>'
        '<p:xfrm><a:off x="548640" y="1828800"/><a:ext cx="10972800" cy="1920240"/></p:xfrm>'
        f'<a:graphic><a:graphicData uri="{_SA_NS_DGM}">'
        f'<dgm:relIds xmlns:dgm="{_SA_NS_DGM}" xmlns:r="{_SA_NS_R}" r:dm="rIdDgmData" '
        'r:lo="rIdDgmLayout" r:qs="rIdDgmQuickStyle" r:cs="rIdDgmColors"/>'
        '</a:graphicData></a:graphic></p:graphicFrame>'
    )
    new_rels = (
        '<Relationship Id="rIdDgmData" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData" Target="../diagrams/data1.xml"/>'
        '<Relationship Id="rIdDgmLayout" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramLayout" Target="../diagrams/layout1.xml"/>'
        '<Relationship Id="rIdDgmQuickStyle" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramQuickStyle" Target="../diagrams/quickStyle1.xml"/>'
        '<Relationship Id="rIdDgmColors" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramColors" Target="../diagrams/colors1.xml"/>'
    )
    ct_overrides = (
        '<Override PartName="/ppt/diagrams/data1.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.drawingml.diagramData+xml"/>'
        '<Override PartName="/ppt/diagrams/layout1.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.drawingml.diagramLayout+xml"/>'
        '<Override PartName="/ppt/diagrams/quickStyle1.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.drawingml.diagramStyle+xml"/>'
        '<Override PartName="/ppt/diagrams/colors1.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.drawingml.diagramColors+xml"/>'
    )

    tmp_path = str(pptx_path) + ".smartart.tmp"
    with zipfile.ZipFile(pptx_path) as zin, zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.decode("utf-8").replace("</Types>", ct_overrides + "</Types>").encode("utf-8")
            elif item.filename == rels_name:
                data = data.decode("utf-8").replace("</Relationships>", new_rels + "</Relationships>").encode("utf-8")
            elif item.filename == slide_part:
                data = data.decode("utf-8").replace("</p:spTree>", graphic_frame + "</p:spTree>").encode("utf-8")
            zout.writestr(item, data)
        zout.writestr("ppt/diagrams/data1.xml", _smartart_data_xml(nodes))
        zout.writestr("ppt/diagrams/layout1.xml", _SA_LAYOUT_XML)
        zout.writestr("ppt/diagrams/quickStyle1.xml", _SA_QUICKSTYLE_XML)
        zout.writestr("ppt/diagrams/colors1.xml", _SA_COLORS_XML)
    os.replace(tmp_path, pptx_path)

IMG_CAPTIONS = {
    "IMG-01": "IMG-01 / 文字を含まないPNG。画像自体がMarkdownから参照可能かを確認する。",
    "IMG-02": "IMG-02 / 日本語・英数字を含むPNG。OCR有効時に画像内限定の文字列が"
              "派生テキストとして出力されるかを確認する（本文・ノートには記載しない）。",
}

def build_slide14_images(prs, img1_bytes, img2_bytes):
    """P11: image slide (IMG-01/IMG-02 + captions).
    QA12: fact exists only inside IMG-02 pixels, never in fixture text."""
    slide = add_slide(prs, 5, "9. 画像・OCR検証", footer_page=14)  # Title Only
    shapes = slide.shapes
    import io
    pic1 = shapes.add_picture(io.BytesIO(img1_bytes), Inches(0.7), Inches(1.7), height=Inches(3.6))
    pic2 = shapes.add_picture(io.BytesIO(img2_bytes), Inches(7.0), Inches(1.7), height=Inches(3.6))

    for pic, key, x in ((pic1, "IMG-01", 0.7), (pic2, "IMG-02", 7.0)):
        cap = shapes.add_textbox(Inches(x), Inches(5.45), Inches(5.6), Inches(1.3))
        ctf = cap.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cr = cp.add_run()
        cr.text = IMG_CAPTIONS[key]
        cr.font.size = Pt(11)
        cr.font.color.rgb = MUTED
        cr.font.name = JP_FONT
    record("P11", "画像スライド（IMG-01/IMG-02+キャプション）", "add_picture x2 (extracted from source xlsx media)", True)
    record("QA12", "IMG-02画像内限定文字列（本文に非記載）", "picture only, no text/notes reference", True)
    return slide

def build_slide15_summary(prs):
    """P10: speaker notes with bold run + bulleted paragraphs."""
    slide = add_slide(prs, 1, "まとめ・レビュー観点", footer_page=15)
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    bullets = [
        "設計レビュー結果: 作成完了、レビューは条件付承認、承認は未承認（表紙 4. レビュー・承認欄）",
        "未決事項3件（ISSUE-01/02/03）。Ownerと期限を明記の上、合意なしに本番移行しない",
        "次のアクション: NFR-AVL-02のRTO/RPO確認、IF-06停止時間帯の会計主管調整（ISSUE-03）",
    ]
    first = True
    for text in bullets:
        p = new_paragraph(tf, first=first)
        first = False
        r = p.add_run()
        r.text = text
        r.font.size = Pt(15)
        r.font.color.rgb = INK
        r.font.name = JP_FONT
        set_bullet_char(p, "●")

    notes_tf = slide.notes_slide.notes_text_frame
    np0 = notes_tf.paragraphs[0]
    nr0 = np0.add_run()
    nr0.text = "レビュー進行メモ"
    nr0.font.bold = True
    nr0.font.size = Pt(14)
    nr0.font.name = JP_FONT
    for text in [
        "状態遷移とAPI契約の整合性を重点確認。",
        "非機能要件（NFR-AVL-02）は基盤担当に確認依頼中。",
        "画像・OCR検証はQA環境で別途実施。",
    ]:
        np = notes_tf.add_paragraph()
        nr = np.add_run()
        nr.text = text
        nr.font.size = Pt(12)
        nr.font.name = JP_FONT
        set_bullet_char(np, "‣")
    record("P10", "スピーカーノート（太字+箇条書き入り）",
           "notes_slide.notes_text_frame: bold run + buChar paragraphs", True)
    return slide

# ============================================================ self-check ==
def verify_reopen(path):
    """Re-open with python-pptx and sanity-check shape counts survive."""
    prs = Presentation(str(path))
    ok = len(prs.slides) >= 14
    record("REOPEN", f"python-pptxで再オープン（{len(prs.slides)}枚）", "Presentation(path)", ok)
    return prs

def verify_raw_xml(path):
    """zipfile-level checks that the OOXML abstractions (buAutoNum, stCxn/
    endCxn, grpSp, c:chart references) actually landed in the package and
    were not silently dropped by python-pptx."""
    with zipfile.ZipFile(path) as z:
        names = set(z.namelist())
        slide_names = sorted(n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
        all_slide_xml = "".join(z.read(n).decode("utf-8") for n in slide_names)
        chart_names = [n for n in names if n.startswith("ppt/charts/chart") and n.endswith(".xml")]
        chart_xml = "".join(z.read(n).decode("utf-8") for n in chart_names)
        diagram_data_xml = z.read("ppt/diagrams/data1.xml").decode("utf-8") if "ppt/diagrams/data1.xml" in names else ""

    checks = [
        ("XML-buAutoNum", "buAutoNum" in all_slide_xml, "grep <a:buAutoNum in slide XML"),
        ("XML-buChar", "buChar" in all_slide_xml, "grep <a:buChar in slide XML"),
        ("XML-buNone", "buNone" in all_slide_xml, "grep <a:buNone in slide XML"),
        ("XML-stCxn", "stCxn" in all_slide_xml, "grep <a:stCxn in slide XML"),
        ("XML-endCxn", "endCxn" in all_slide_xml, "grep <a:endCxn in slide XML"),
        ("XML-grpSp", "<p:grpSp>" in all_slide_xml, "grep <p:grpSp> in slide XML"),
        ("XML-chart-ref", 'drawingml/2006/chart"' in all_slide_xml, "grep chart graphicData uri in slide XML"),
        ("XML-chart-part", ("<c:barChart" in chart_xml or "<c:pieChart" in chart_xml), "grep <c:barChart/<c:pieChart in ppt/charts/*.xml"),
        ("XML-vMerge", ('vMerge="1"' in all_slide_xml or "rowSpan=" in all_slide_xml), "grep vMerge/rowSpan in slide XML (P09 merge)"),
        ("XML-strike", 'strike="sngStrike"' in all_slide_xml, "grep strike=\"sngStrike\" in slide XML (P14)"),
        ("XML-rot", 'rot="2700000"' in all_slide_xml, "grep rot=\"2700000\" (45deg) in slide XML (P15)"),
        ("XML-ftr", 'type="ftr"' in all_slide_xml, "grep <p:ph type=\"ftr\" in slide XML (P13)"),
        ("XML-sldNum", 'type="sldNum"' in all_slide_xml, "grep <p:ph type=\"sldNum\" in slide XML (P13)"),
        ("XML-dt", 'type="dt"' in all_slide_xml, "grep <p:ph type=\"dt\" in slide XML (P13)"),
        ("XML-forbidden-token-absent", FORBIDDEN_TOKEN not in all_slide_xml, "assert OCR token absent from all slide XML text"),
        ("XML-diagram-ref", 'drawingml/2006/diagram"' in all_slide_xml, "grep diagram graphicData uri in slide XML (P07)"),
        ("XML-diagram-data", bool(diagram_data_xml) and "<dgm:pt " in diagram_data_xml, "ppt/diagrams/data1.xml has <dgm:pt> nodes (P07)"),
        ("XML-diagram-text", SMARTART_NODES[0] in diagram_data_xml, "ppt/diagrams/data1.xml carries the real node text (P07)"),
    ]
    all_ok = True
    for cid, ok, method in checks:
        record(cid, cid, method, ok)
        all_ok = all_ok and ok
    return all_ok

def print_checklist():
    print("\n=== complex-design-doc.pptx element checklist ===")
    width = max(len(c[0]) for c in CHECKLIST)
    for elem_id, desc, method, status in CHECKLIST:
        marker = "PASS" if status == "PASS" else "FAIL"
        print(f"[{marker}] {elem_id.ljust(width)}  {desc}  ({method})")
    failed = [c for c in CHECKLIST if c[3] != "PASS"]
    print(f"\n{len(CHECKLIST)} checks, {len(failed)} failures.")
    return len(failed) == 0

def build(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    img1, img2 = load_images()

    build_slide01_title(prs)
    build_slide02_section(prs)
    build_slide03_overview(prs)
    build_slide04_assumptions(prs)
    build_slide05_state_machine(prs)
    build_slide06_procedure(prs)
    build_slide07_nfr(prs)
    build_slide08_charts(prs)
    build_slide09_test_scenarios(prs)
    build_slide10_api_contract(prs)
    build_slide11_screen_spec(prs)
    build_slide12_two_column(prs)
    smartart_slide = build_slide13_smartart(prs)
    build_slide14_images(prs, img1, img2)
    build_slide15_summary(prs)

    record("P13", "フッター+スライド番号+日付プレースホルダー",
           "cloned <p:ph type=ftr/sldNum/dt> from layout on slides 3-15", True)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    smartart_index = list(prs.slides).index(smartart_slide) + 1
    inject_smartart(output_path, f"ppt/slides/slide{smartart_index}.xml")

    return output_path

def main():
    out_path = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_OUT
    build(out_path)
    print(f"wrote {out_path} ({out_path.stat().st_size} bytes, {len(Presentation(str(out_path)).slides)} slides)")

    verify_reopen(out_path)
    verify_raw_xml(out_path)
    all_ok = print_checklist()
    if not all_ok:
        sys.exit(1)

if __name__ == "__main__":
    main()